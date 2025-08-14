import streamlit as st
from pptx import Presentation
import os
from datetime import datetime
import io
import tempfile

# --- Page config ---
st.set_page_config(page_title="📝 PPT Toggle Forms", layout="centered")
st.markdown(
    """
    <style>
    .stApp {background-color: #a8d0e6;}
    .block-container {background-color: #fff; padding: 2rem; border-radius: 10px;}
    h1 {color: darkblue;}
    .streamlit-expanderHeader {color: black !important;}
    </style>
    """,
    unsafe_allow_html=True
)
st.markdown("<h1 style='text-align: center;'>📊 PPT Data Entry Toggle Forms</h1>", unsafe_allow_html=True)
st.write("---")

# --- Initialize session state ---
if "Case_Study_Template_1" not in st.session_state:
    st.session_state.Case_Study_Template_1 = False
if "Case_Study_Template_2" not in st.session_state:
    st.session_state.Case_Study_Template_2 = False

# --- Buttons to toggle forms ---
col1, col2 = st.columns(2)
with col1:
    if st.button("Case Study For Template-1"):
        st.session_state.Case_Study_Template_1 = True
        st.session_state.Case_Study_Template_2 = False
with col2:
    if st.button("Case Study For Template-2"):
        st.session_state.Case_Study_Template_2 = True
        st.session_state.Case_Study_Template_1 = False

st.write("---")

date_format = "%d-%m-%Y"

# --- Helper function to generate PPT ---
def generate_ppt(template_path, replacements, images):
    prs = Presentation(template_path)

    def replace_placeholder(para, replacements):
        full_text = "".join(run.text for run in para.runs)
        for key, value in replacements.items():
            full_text = full_text.replace(key, value)
        for run in para.runs:
            run.text = ""
        if para.runs:
            para.runs[0].text = full_text
        else:
            para.add_run().text = full_text

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                replaced_image = False
                for key, file in images.items():
                    if key in shape.text and file:
                        shape.text = ""
                        left, top, width, height = shape.left, shape.top, shape.width, shape.height
                        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                            tmp.write(file.getvalue())
                            tmp_path = tmp.name
                        slide.shapes.add_picture(tmp_path, left, top, width, height)
                        os.unlink(tmp_path)
                        replaced_image = True
                        break
                if not replaced_image:
                    for para in shape.text_frame.paragraphs:
                        replace_placeholder(para, replacements)

    ppt_bytes = io.BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes

# --- Form 1: Case Study Template-1 ---
ppt1_bytes = None
if st.session_state.Case_Study_Template_1:
    with st.form("form1"):
        st.header("📝 Case Study For Template-1")
        plant_name = st.text_input("Plant Name")
        equipment_name = st.text_input("Equipment")
        case_enabler = st.text_input("Case Enabler")
        downtime_hours = st.text_input("Downtime Hours")
        observation_date = st.date_input("Observation Date")
        observation = st.text_input("Observation")
        date_recommendation = st.date_input("Date of Recommendation")
        recommendation = st.text_input("Recommendation")
        date_corrective_action = st.date_input("Date of Corrective Action Taken")
        corrective_action_details = st.text_input("Corrective Action Details")
        date_closed_report = st.date_input("Date of Closed Report")
        closed_report_status = st.text_input("Closed Report Status")
        machine_details = st.text_area("Machine Details")
        trend_for_image1 = st.text_input("Trend for Image-1")
        trend_for_image2 = st.text_input("Trend for Image-2")
        equipment_image = st.file_uploader("📷 Upload Equipment Image", type=['png','jpg','jpeg'])
        trend_image1 = st.file_uploader("📷 Upload Trend Image 1", type=['png','jpg','jpeg'])
        trend_image2 = st.file_uploader("📷 Upload Trend Image 2", type=['png','jpg','jpeg'])
        submit1 = st.form_submit_button("✅ Generate PPT Template-1")

    if submit1:
        template_path = "template1.pptx"
        replacements = {
            "{{Plant Name}}": plant_name,
            "{{Equipment}}": equipment_name,
            "{{Case Enabler}}": case_enabler,
            "{{Downtime Hours}}": downtime_hours,
            "{{Observation Date}}": observation_date.strftime(date_format),
            "{{Observation}}": observation,
            "{{Date of Recommendation}}": date_recommendation.strftime(date_format),
            "{{Recommendation}}": recommendation,
            "{{Date of Corrective action Taken}}": date_corrective_action.strftime(date_format),
            "{{Corrective Action Details}}": corrective_action_details,
            "{{Date of closed Report}}": date_closed_report.strftime(date_format),
            "{{Closed Report Status}}": closed_report_status,
            "{{Machine Details}}": machine_details,
            "{{Trend for Image-1}}": trend_for_image1,
            "{{Trend for Image-2}}": trend_for_image2
        }
        images = {
            "{{Equipment Image}}": equipment_image,
            "{{Trend Image 1}}": trend_image1,
            "{{Trend Image 2}}": trend_image2
        }
        ppt1_bytes = generate_ppt(template_path, replacements, images)
        st.success("🎉 PPT Template-1 generated successfully!")

if ppt1_bytes:
    file_name = f"{equipment_name.replace(' ','_')}_{plant_name.replace(' ','_')}.pptx"
    st.download_button("⬇️ Download PPT Template-1", ppt1_bytes, file_name=file_name,
                       mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

# --- Form 2: Case Study Template-2 ---
ppt2_bytes = None
if st.session_state.Case_Study_Template_2:
    with st.form("form2"):
        st.header("📝 Case Study For Template-2")
        plant_name2 = st.text_input("Plant Name")
        equipment_name2 = st.text_input("Equipment")
        case_enabler2 = st.text_input("Case Enabler")
        downtime_hours2 = st.text_input("Downtime Hours")
        observation_date2 = st.date_input("Observation Date")
        observation2 = st.text_area("Observation Details")
        date_recommendation2 = st.date_input("Date of Recommendation")
        recommendation2 = st.text_input("Recommendation")
        date_corrective_action2 = st.date_input("Date of Corrective Action Taken")
        corrective_action_details2 = st.text_input("Corrective Action Details")
        date_closed_report2 = st.date_input("Date of Closed Report")
        closed_report_status2 = st.text_input("Closed Report Status")
        machine_details2 = st.text_area("Machine Details")
        trend_for_image1_2 = st.text_input("Trend for Image-1")
        equipment_image2 = st.file_uploader("📷 Upload Equipment Image", type=['png','jpg','jpeg'])
        trend_image1_2 = st.file_uploader("📷 Upload Trend Image 1", type=['png','jpg','jpeg'])
        submit2 = st.form_submit_button("✅ Generate PPT Template-2")

    if submit2:
        template_path = "template2.pptx"
        replacements = {
            "{{Plant Name}}": plant_name2,
            "{{Equipment}}": equipment_name2,
            "{{Case Enabler}}": case_enabler2,
            "{{Downtime Hours}}": downtime_hours2,
            "{{Observation Date}}": observation_date2.strftime(date_format),
            "{{Observation}}": observation2,
            "{{Date of Recommendation}}": date_recommendation2.strftime(date_format),
            "{{Recommendation}}": recommendation2,
            "{{Date of Corrective action Taken}}": date_corrective_action2.strftime(date_format),
            "{{Corrective Action Details}}": corrective_action_details2,
            "{{Date of closed Report}}": date_closed_report2.strftime(date_format),
            "{{Closed Report Status}}": closed_report_status2,
            "{{Machine Details}}": machine_details2,
            "{{Trend for Image-1}}": trend_for_image1_2,
        }
        images = {
            "{{Equipment Image}}": equipment_image2,
            "{{Trend Image 1}}": trend_image1_2
        }
        ppt2_bytes = generate_ppt(template_path, replacements, images)
        st.success("🎉 PPT Template-2 generated successfully!")

if ppt2_bytes:
    file_name = f"{equipment_name2.replace(' ','_')}_{plant_name2.replace(' ','_')}.pptx"
    st.download_button("⬇️ Download PPT Template-2", ppt2_bytes, file_name=file_name,
                       mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
